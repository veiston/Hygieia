import os
from pptx import Presentation
from docx import Document
import PyPDF2

def anyReader(file_path=None):
    # quick file reader for PDF, PPTX, DOCX
    if not file_path:
        return ''
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        try:
            with open(file_path, 'rb') as f:
                pdf = PyPDF2.PdfReader(f)
                out = ''
                for i, page in enumerate(pdf.pages):
                    t = page.extract_text()
                    if t:
                        out += f"\n\nPage {i+1}\n{t}"
                return out.strip()
        except Exception as e:
            return f'PDF read error: {e}'
    if ext == '.pptx':
        try:
            prs = Presentation(file_path)
            runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    tf = getattr(shape, "text_frame", None)
                    if tf is None:
                        continue
                    for p in tf.paragraphs:
                        for r in p.runs:
                            runs.append(r.text)
                        runs.append('\n')
            return '\n'.join(runs).strip()
        except Exception as e:
            return f'PPTX read error: {e}'
        
    if ext == '.docx':
        try:
            doc = Document(file_path)
            return '\n'.join([p.text for p in doc.paragraphs if p.text.strip()]).strip()
        except Exception as e:
            return f'DOCX read error: {e}'
    return 'Unsupported filetype'

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        print(anyReader(sys.argv[1]))
    else:
        print('Usage: python anyFileRead.py <file_path>')